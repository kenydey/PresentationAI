import asyncio
import os
from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxChartBoxModel,
    PptxChartSeriesModel,
    PptxFillModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxSlideModel,
    PptxTableBoxModel,
)
from services.pptx_presentation_creator import PptxPresentationCreator
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from utils.slide_to_pptx_converter import convert_presentation_to_pptx_model


pptx_model = PptxPresentationModel(
    slides=[
        PptxSlideModel(
            shapes=[
                PptxAutoShapeBoxModel(
                    type=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    position=PptxPositionModel(
                        left=20,
                        right=20,
                        width=100,
                        height=100,
                    ),
                    fill=PptxFillModel(
                        color="000000",
                        opacity=0.5,
                    ),
                )
            ]
        )
    ]
)


def test_pptx_creator():
    temp_dir = "/tmp/presenton"
    pptx_creator = PptxPresentationCreator(pptx_model, temp_dir)
    asyncio.run(pptx_creator.create_ppt())
    os.makedirs("debug", exist_ok=True)
    pptx_creator.save("debug/test.pptx")


def test_pptx_table_chart_export():
    """验证表格与图表可正确导出为原生 PPTX。"""
    slides = [
        {
            "content": {
                "title": "季度数据",
                "table": {
                    "headers": ["指标", "Q1", "Q2"],
                    "rows": [["收入", "120", "150"], ["成本", "80", "95"]],
                },
            },
        },
        {
            "content": {
                "title": "趋势",
                "chart": {
                    "type": "bar",
                    "categories": ["Q1", "Q2", "Q3"],
                    "series": [{"name": "产品A", "data": [100, 120, 140]}],
                },
            },
        },
    ]
    model = convert_presentation_to_pptx_model(slides, title="表格图表测试")
    assert len(model.slides) == 2
    table_shapes = [s for s in model.slides[0].shapes if isinstance(s, PptxTableBoxModel)]
    chart_shapes = [s for s in model.slides[1].shapes if isinstance(s, PptxChartBoxModel)]
    assert len(table_shapes) == 1
    assert len(chart_shapes) == 1
    assert table_shapes[0].headers == ["指标", "Q1", "Q2"]
    assert chart_shapes[0].chart_type == "bar"

    temp_dir = "/tmp/presenton"
    creator = PptxPresentationCreator(model, temp_dir)
    asyncio.run(creator.create_ppt())
    os.makedirs("debug", exist_ok=True)
    creator.save("debug/test_table_chart.pptx")
    assert os.path.exists("debug/test_table_chart.pptx")
