"""
文本自动适应（Text Autofit）：弥补 python-pptx 缺乏渲染引擎的短板。
利用 Pillow 的字体矩阵模拟计算文本尺寸，在服务端预判溢出并调整字号。
"""

from typing import Optional

# python-pptx 常量
try:
    from pptx.enum.text import MSO_AUTO_SIZE
except ImportError:
    MSO_AUTO_SIZE = None


def apply_text_frame_autofit_properties(text_frame) -> None:
    """
    为文本框设置 word_wrap 和 auto_size，促使 PowerPoint 打开时自动重排。
    策略 2：枚举指令注入
    """
    if text_frame is None:
        return
    text_frame.word_wrap = True
    if MSO_AUTO_SIZE:
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        except Exception:
            pass


def compute_text_fitting_font_size(
    text: str,
    font_path: str,
    max_width_pt: float,
    max_height_pt: float,
    start_pt: int = 24,
    min_pt: int = 10,
) -> Optional[int]:
    """
    使用 Pillow 模拟文本尺寸，计算能适应给定矩形的最小字号。
    策略 3：基于 Pillow 的字体矩阵模拟运算。
    返回适应的字号，若无法适应则返回 min_pt。
    """
    try:
        from PIL import ImageFont, ImageDraw
    except ImportError:
        return None

    # pt 转 px 近似：1pt ≈ 1.33px @ 96dpi
    scale = 1.33
    max_w_px = int(max_width_pt * scale)
    max_h_px = int(max_height_pt * scale)

    for pt in range(start_pt, min_pt - 1, -1):
        try:
            font = ImageFont.truetype(font_path, pt)
        except OSError:
            font = ImageFont.load_default()
        bbox = ImageDraw.Draw(Image.new("RGB", (1, 1))).textbbox(
            (0, 0), text, font=font
        )
        w = bbox[2] - bbox[0]
        h = bbox[3] - bbox[1]
        if w <= max_w_px and h <= max_h_px:
            return pt
    return min_pt
